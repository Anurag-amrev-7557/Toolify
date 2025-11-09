from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os
import io
import requests

app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "*", "methods": ["GET", "POST", "OPTIONS"], "allow_headers": "*"}})

@app.after_request
def after_request(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    # Allow the client to use custom headers (e.g. X-Temp-Upload) by echoing the
    # requested headers from the preflight or falling back to a sensible default.
    requested = request.headers.get('Access-Control-Request-Headers')
    if requested:
        response.headers.add('Access-Control-Allow-Headers', requested)
    else:
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type, X-Temp-Upload, Authorization')
    response.headers.add('Access-Control-Allow-Methods', 'GET,POST,OPTIONS')
    return response

UPLOAD_FOLDER = 'temp'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

@app.route('/api/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok'})

@app.route('/api/process/<tool>', methods=['POST', 'OPTIONS'])
def process_tool(tool):
    if request.method == 'OPTIONS':
        return '', 204
    try:
        if tool == 'word-counter':
            file = request.files.get('file')
            text = file.read().decode('utf-8') if file else request.form.get('text', '')
            words = len(text.split())
            chars = len(text)
            lines = len(text.splitlines())
            return jsonify({'words': words, 'characters': chars, 'lines': lines})
        
        elif tool == 'json-formatter':
            file = request.files.get('file')
            text = file.read().decode('utf-8') if file else request.form.get('text', '')
            import json
            formatted = json.dumps(json.loads(text), indent=2)
            return jsonify({'formatted': formatted})
        
        elif tool == 'base64-encoder':
            import base64
            text = request.form.get('text', '')
            action = request.form.get('action', 'encode')
            if action == 'encode':
                result = base64.b64encode(text.encode()).decode()
            else:
                result = base64.b64decode(text.encode()).decode()
            return jsonify({'result': result})
        
        elif tool == 'uuid-generator':
            import uuid
            return jsonify({'uuid': str(uuid.uuid4())})
        
        elif tool == 'password-generator':
            import random, string
            length = int(request.form.get('length', 16))
            chars = string.ascii_letters + string.digits + string.punctuation
            password = ''.join(random.choice(chars) for _ in range(length))
            return jsonify({'password': password})
        
        elif tool == 'hash-generator':
            import hashlib
            text = request.form.get('text', '')
            return jsonify({
                'md5': hashlib.md5(text.encode()).hexdigest(),
                'sha256': hashlib.sha256(text.encode()).hexdigest()
            })
        
        elif tool == 'case-converter':
            text = request.form.get('text', '')
            case_type = request.form.get('type', 'upper')
            result = text.upper() if case_type == 'upper' else text.lower() if case_type == 'lower' else text.title()
            return jsonify({'result': result})
        
        elif tool == 'pdf-merger':
            from pypdf import PdfMerger
            files = request.files.getlist('files')
            merger = PdfMerger()
            for f in files:
                merger.append(f)
            output = os.path.join(UPLOAD_FOLDER, 'merged.pdf')
            merger.write(output)
            merger.close()
            return send_file(output, as_attachment=True)
        
        elif tool == 'image-compressor':
            from PIL import Image
            
            file = request.files.get('file')
            quality = int(request.form.get('quality', 80))
            format_choice = request.form.get('format', 'original')
            max_width = int(request.form.get('maxWidth', 0))
            max_height = int(request.form.get('maxHeight', 0))
            maintain_aspect = request.form.get('maintainAspect', 'true') == 'true'
            rotation = int(request.form.get('rotation', 0))
            flip_h = request.form.get('flipH', 'false') == 'true'
            flip_v = request.form.get('flipV', 'false') == 'true'
            brightness = float(request.form.get('brightness', 100)) / 100
            contrast = float(request.form.get('contrast', 100)) / 100
            
            original_filename = file.filename
            file_ext = original_filename.rsplit('.', 1)[-1].lower()
            
            img = Image.open(file)
            original_size = img.size
            
            # Apply rotation
            if rotation != 0:
                img = img.rotate(-rotation, expand=True)
            
            # Apply flips
            if flip_h:
                img = img.transpose(Image.FLIP_LEFT_RIGHT)
            if flip_v:
                img = img.transpose(Image.FLIP_TOP_BOTTOM)
            
            # Apply brightness and contrast
            from PIL import ImageEnhance
            if brightness != 1.0:
                enhancer = ImageEnhance.Brightness(img)
                img = enhancer.enhance(brightness)
            if contrast != 1.0:
                enhancer = ImageEnhance.Contrast(img)
                img = enhancer.enhance(contrast)
            
            # Resize if dimensions specified with smart resampling
            if max_width > 0 or max_height > 0:
                original_width, original_height = img.size
                if maintain_aspect:
                    img.thumbnail((max_width or 999999, max_height or 999999), Image.Resampling.LANCZOS)
                else:
                    new_width = max_width if max_width > 0 else img.width
                    new_height = max_height if max_height > 0 else img.height
                    img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                
                # Apply sharpening after resize to maintain perceived quality
                from PIL import ImageFilter
                if img.size[0] < original_width * 0.8:  # Only if significantly downsized
                    img = img.filter(ImageFilter.UnsharpMask(radius=1, percent=100, threshold=3))
            
            # Determine output format
            if format_choice == 'original':
                output_ext = file_ext
            else:
                output_ext = format_choice
            
            # Handle format conversion
            if output_ext in ['jpg', 'jpeg']:
                if img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                    img = background
                save_format = 'JPEG'
            elif output_ext == 'png':
                save_format = 'PNG'
            elif output_ext == 'webp':
                save_format = 'WEBP'
            else:
                save_format = 'JPEG'
                output_ext = 'jpg'
            
            # Apply smart pre-compression optimization
            from PIL import ImageFilter
            
            # Slight noise reduction for better compression (especially for photos)
            if quality < 85 and img.mode in ('RGB', 'RGBA'):
                img = img.filter(ImageFilter.MedianFilter(size=3))
            
            # Save with appropriate settings for maximum quality retention
            output_buffer = io.BytesIO()
            
            # Prepare EXIF data stripping for smaller size
            save_kwargs = {}
            
            if save_format == 'JPEG':
                # Use subsampling for better quality at high compression
                if quality >= 90:
                    subsampling = 0  # 4:4:4 (best quality)
                elif quality >= 80:
                    subsampling = 1  # 4:2:2
                else:
                    subsampling = 2  # 4:2:0 (default)
                
                # Strip metadata for smaller size
                save_kwargs = {
                    'format': save_format,
                    'quality': quality,
                    'optimize': True,
                    'progressive': True,
                    'subsampling': subsampling,
                    'qtables': 'web_high' if quality >= 85 else 'web_low',
                    'exif': b''  # Remove EXIF data
                }
                
            elif save_format == 'PNG':
                # Use adaptive filtering and optimize for better compression
                compress_level = int(9 - (quality / 100 * 9))
                
                # Smart palette conversion for smaller size
                if img.mode == 'RGB' and quality < 90:
                    # Use quantization for better color selection
                    img = img.quantize(colors=256, method=2, dither=1)
                elif img.mode == 'RGBA' and quality < 90:
                    # Preserve alpha channel
                    alpha = img.split()[-1]
                    rgb = img.convert('RGB').quantize(colors=256, method=2, dither=1)
                    rgb.putalpha(alpha)
                    img = rgb
                
                save_kwargs = {
                    'format': save_format,
                    'optimize': True,
                    'compress_level': compress_level
                }
                
            elif save_format == 'WEBP':
                # WebP with advanced options
                if quality >= 95:
                    save_kwargs = {
                        'format': save_format,
                        'lossless': True,
                        'quality': 100,
                        'method': 6
                    }
                else:
                    save_kwargs = {
                        'format': save_format,
                        'quality': quality,
                        'method': 6,  # Slowest but best compression
                        'exact': quality >= 85,  # Preserve exact colors for high quality
                        'minimize_size': True  # Additional size optimization
                    }
            
            img.save(output_buffer, **save_kwargs)
            
            output_buffer.seek(0)
            compressed_data = output_buffer.getvalue()
            
            # Additional optimization: Try mozjpeg-like optimization for JPEG
            if save_format == 'JPEG' and quality < 95:
                # Re-compress with slightly different settings for better compression
                temp_buffer = io.BytesIO()
                temp_img = Image.open(io.BytesIO(compressed_data))
                temp_img.save(
                    temp_buffer,
                    format='JPEG',
                    quality=quality,
                    optimize=True,
                    progressive=True
                )
                temp_data = temp_buffer.getvalue()
                # Use smaller result
                if len(temp_data) < len(compressed_data):
                    compressed_data = temp_data
            
            output_filename = f'compressed.{output_ext}'
            output_path = os.path.join(UPLOAD_FOLDER, output_filename)
            
            with open(output_path, 'wb') as f:
                f.write(compressed_data)
            
            return send_file(output_path, as_attachment=True, download_name=output_filename)
        
        
        elif tool == 'pdf-splitter':
            from pypdf import PdfReader, PdfWriter
            import zipfile
            file = request.files.get('file')
            mode = request.form.get('mode', 'all')
            reader = PdfReader(file)
            total_pages = len(reader.pages)
            # Quick page count endpoint: return only the number of pages when requested
            if mode == 'count':
                return jsonify({'pages': total_pages})
            if mode == 'all':
                zip_path = os.path.join(UPLOAD_FOLDER, 'split-pages.zip')
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    for i in range(total_pages):
                        writer = PdfWriter()
                        writer.add_page(reader.pages[i])
                        pdf_path = os.path.join(UPLOAD_FOLDER, f'page_{i+1}.pdf')
                        with open(pdf_path, 'wb') as f:
                            writer.write(f)
                        zipf.write(pdf_path, f'page_{i+1}.pdf')
                        os.remove(pdf_path)
                return send_file(zip_path, as_attachment=True)
            elif mode == 'range':
                try:
                    start = int(request.form.get('start', 1))
                    end = int(request.form.get('end', total_pages))
                except Exception:
                    return jsonify({'error': 'Invalid start or end values'}), 400

                # Validate ranges
                if start < 1 or end < 1 or start > end:
                    return jsonify({'error': 'Invalid page range: ensure 1 <= start <= end'}), 400
                if start > total_pages:
                    return jsonify({'error': 'Start page is greater than total pages'}), 400
                # cap the end to total_pages
                end = min(end, total_pages)

                writer = PdfWriter()
                for i in range(start - 1, end):
                    writer.add_page(reader.pages[i])
                output = os.path.join(UPLOAD_FOLDER, 'split-range.pdf')
                with open(output, 'wb') as f:
                    writer.write(f)
                return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-compress':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            temp_input = os.path.join(UPLOAD_FOLDER, 'temp_input.pdf')
            output = os.path.join(UPLOAD_FOLDER, 'compressed.pdf')
            file.save(temp_input)
            try:
                reader = PdfReader(temp_input)
                writer = PdfWriter()
                for page in reader.pages:
                    page.compress_content_streams()
                    writer.add_page(page)
                with open(output, 'wb') as f:
                    writer.write(f)
                return send_file(output, as_attachment=True, download_name='compressed.pdf')
            finally:
                if os.path.exists(temp_input):
                    os.remove(temp_input)
        
        elif tool == 'pdf-to-jpg':
            import zipfile
            import fitz
            file = request.files.get('file')
            temp_input = os.path.join(UPLOAD_FOLDER, 'temp_input.pdf')
            file.save(temp_input)
            try:
                doc = fitz.open(temp_input)
                zip_path = os.path.join(UPLOAD_FOLDER, 'pdf-images.zip')
                with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        pix = page.get_pixmap(dpi=200)
                        img_path = os.path.join(UPLOAD_FOLDER, f'page_{page_num+1}.jpg')
                        pix.save(img_path)
                        zipf.write(img_path, f'page_{page_num+1}.jpg')
                        os.remove(img_path)
                doc.close()
                return send_file(zip_path, as_attachment=True, mimetype='application/zip')
            finally:
                if os.path.exists(temp_input):
                    os.remove(temp_input)
        
        elif tool == 'jpg-to-pdf':
            from PIL import Image
            from pypdf import PdfWriter
            # Support three modes for jpg->pdf:
            # 1) Temporary single-file upload (X-Temp-Upload header) -> save file and return a fileId
            # 2) Assemble from previously uploaded fileIds (form field 'fileIds') -> build PDF from saved files
            # 3) Direct combined upload (files field) -> immediate assemble (legacy behavior)

            # Mode 1: temporary upload (used by client to upload individual files with progress)
            if request.headers.get('X-Temp-Upload') or request.form.get('tempUpload'):
                file = request.files.get('file')
                if not file:
                    return jsonify({'error': 'No file uploaded'}), 400
                # generate unique filename to store
                import uuid
                orig = file.filename or 'upload'
                ext = os.path.splitext(orig)[1] or '.jpg'
                file_id = f"{uuid.uuid4().hex}{ext}"
                save_path = os.path.join(UPLOAD_FOLDER, file_id)
                file.save(save_path)
                return jsonify({'fileId': file_id})

            # Mode 2: assemble from fileIds
            file_ids = request.form.getlist('fileIds')
            if file_ids:
                images = []
                for fid in file_ids:
                    path = os.path.join(UPLOAD_FOLDER, fid)
                    if not os.path.exists(path):
                        return jsonify({'error': f'File not found: {fid}'}), 400
                    img = Image.open(path)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    images.append(img)
                if not images:
                    return jsonify({'error': 'No images provided'}), 400
                output = os.path.join(UPLOAD_FOLDER, 'images-to-pdf.pdf')
                images[0].save(output, save_all=True, append_images=images[1:]) if len(images) > 1 else images[0].save(output)
                return send_file(output, as_attachment=True)

            # Mode 3: legacy direct upload with multiple files
            files = request.files.getlist('files')
            if not files:
                return jsonify({'error': 'No files uploaded'}), 400
            images = []
            for f in files:
                img = Image.open(f)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                images.append(img)
            output = os.path.join(UPLOAD_FOLDER, 'images-to-pdf.pdf')
            images[0].save(output, save_all=True, append_images=images[1:]) if len(images) > 1 else images[0].save(output)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-rotate':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            angle = int(request.form.get('angle', 90))
            reader = PdfReader(file)
            writer = PdfWriter()
            for page in reader.pages:
                page.rotate(angle)
                writer.add_page(page)
            output = os.path.join(UPLOAD_FOLDER, 'rotated.pdf')
            with open(output, 'wb') as f:
                writer.write(f)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-watermark':
            from pypdf import PdfReader, PdfWriter, PageObject
            from reportlab.pdfgen import canvas
            file = request.files.get('file')
            text = request.form.get('text', 'WATERMARK')
            temp_input = os.path.join(UPLOAD_FOLDER, 'temp_watermark_input.pdf')
            file.save(temp_input)
            try:
                reader = PdfReader(temp_input)
                writer = PdfWriter()
                
                for page in reader.pages:
                    packet = io.BytesIO()
                    page_width = float(page.mediabox.width)
                    page_height = float(page.mediabox.height)
                    can = canvas.Canvas(packet, pagesize=(page_width, page_height))
                    can.setFillAlpha(0.3)
                    can.setFillColorRGB(0.5, 0.5, 0.5)
                    can.setFont('Helvetica', 50)
                    can.rotate(45)
                    can.drawString(page_width/3, page_height/3, text)
                    can.save()
                    packet.seek(0)
                    watermark = PdfReader(packet)
                    page.merge_page(watermark.pages[0])
                    writer.add_page(page)
                
                output = os.path.join(UPLOAD_FOLDER, 'watermarked.pdf')
                with open(output, 'wb') as f:
                    writer.write(f)
                return send_file(output, as_attachment=True)
            finally:
                if os.path.exists(temp_input):
                    os.remove(temp_input)
        
        elif tool == 'pdf-unlock':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            password = request.form.get('password', '')
            reader = PdfReader(file)
            if reader.is_encrypted:
                reader.decrypt(password)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            output = os.path.join(UPLOAD_FOLDER, 'unlocked.pdf')
            with open(output, 'wb') as f:
                writer.write(f)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-protect':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            password = request.form.get('password', '')
            reader = PdfReader(file)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.encrypt(password)
            output = os.path.join(UPLOAD_FOLDER, 'protected.pdf')
            with open(output, 'wb') as f:
                writer.write(f)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-organize':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            order = request.form.get('order', '')
            reader = PdfReader(file)
            writer = PdfWriter()
            page_order = [int(x)-1 for x in order.split(',')]
            for i in page_order:
                if 0 <= i < len(reader.pages):
                    writer.add_page(reader.pages[i])
            output = os.path.join(UPLOAD_FOLDER, 'organized.pdf')
            with open(output, 'wb') as f:
                writer.write(f)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-page-numbers':
            from pypdf import PdfReader, PdfWriter
            from reportlab.pdfgen import canvas
            from reportlab.lib.colors import HexColor
            file = request.files.get('file')
            if not file:
                return jsonify({'error': 'No file provided'}), 400
            try:
                reader = PdfReader(file)
                writer = PdfWriter()
                position = request.form.get('position', 'bottom-center')
                start_num = int(request.form.get('startNumber', 1))
                font_size = int(request.form.get('fontSize', 12))
                format_str = request.form.get('format', 'Page {n}')
                color_hex = request.form.get('color', '#000000')
                
                for i, page in enumerate(reader.pages):
                    packet = io.BytesIO()
                    page_width = float(page.mediabox.width)
                    page_height = float(page.mediabox.height)
                    can = canvas.Canvas(packet, pagesize=(page_width, page_height))
                    can.setFont('Helvetica', font_size)
                    can.setFillColor(HexColor(color_hex))
                    
                    text = format_str.replace('{n}', str(i + start_num))
                    text_width = can.stringWidth(text, 'Helvetica', font_size)
                    
                    if 'bottom' in position:
                        y = 20
                    elif 'top' in position:
                        y = page_height - 30
                    else:
                        y = page_height / 2
                    
                    if 'center' in position:
                        x = (page_width - text_width) / 2
                    elif 'right' in position:
                        x = page_width - text_width - 30
                    else:
                        x = 30
                    
                    can.drawString(x, y, text)
                    can.save()
                    packet.seek(0)
                    number_pdf = PdfReader(packet)
                    page.merge_page(number_pdf.pages[0])
                    writer.add_page(page)
                output = os.path.join(UPLOAD_FOLDER, 'numbered.pdf')
                with open(output, 'wb') as f:
                    writer.write(f)
                return send_file(output, as_attachment=True)
            except Exception as e:
                return jsonify({'error': str(e)}), 500
        
        elif tool == 'pdf-extract-pages':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            pages = request.form.get('pages', '')
            reader = PdfReader(file)
            writer = PdfWriter()
            page_list = [int(x)-1 for x in pages.split(',')]
            for i in page_list:
                if 0 <= i < len(reader.pages):
                    writer.add_page(reader.pages[i])
            output = os.path.join(UPLOAD_FOLDER, 'extracted.pdf')
            with open(output, 'wb') as f:
                writer.write(f)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-delete-pages':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            pages = request.form.get('pages', '')
            reader = PdfReader(file)
            writer = PdfWriter()
            delete_list = [int(x)-1 for x in pages.split(',')]
            for i, page in enumerate(reader.pages):
                if i not in delete_list:
                    writer.add_page(page)
            output = os.path.join(UPLOAD_FOLDER, 'deleted-pages.pdf')
            with open(output, 'wb') as f:
                writer.write(f)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-repair':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            reader = PdfReader(file, strict=False)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            output = os.path.join(UPLOAD_FOLDER, 'repaired.pdf')
            with open(output, 'wb') as f:
                writer.write(f)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-to-pdfa':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            reader = PdfReader(file)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.add_metadata({'/Title': 'PDF/A Document'})
            output = os.path.join(UPLOAD_FOLDER, 'pdfa.pdf')
            with open(output, 'wb') as f:
                writer.write(f)
            return send_file(output, as_attachment=True)
        
        elif tool == 'html-to-pdf':
            from weasyprint import HTML
            html_content = request.form.get('html', '')
            output = os.path.join(UPLOAD_FOLDER, 'html-to-pdf.pdf')
            HTML(string=html_content).write_pdf(output)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-ocr':
            import base64
            
            file = request.files.get('file')
            if not file:
                return jsonify({'error': 'No file provided'}), 400
            
            try:
                pdf_base64 = base64.b64encode(file.read()).decode()
                
                payload = {
                    'apikey': 'K83701879288957',
                    'base64Image': f'data:application/pdf;base64,{pdf_base64}',
                    'language': 'eng',
                    'isOverlayRequired': False
                }
                
                response = requests.post('https://api.ocr.space/parse/image', data=payload)
                result = response.json()
                
                if result.get('ParsedResults'):
                    text = '\n\n'.join([page['ParsedText'] for page in result['ParsedResults']])
                    return jsonify({'text': text})
                else:
                    return jsonify({'error': result.get('ErrorMessage', 'OCR failed')}), 500
            except Exception as e:
                return jsonify({'error': f'OCR failed: {str(e)}'}), 500
        
        elif tool == 'pdf-sign':
            from pypdf import PdfReader, PdfWriter
            file = request.files.get('file')
            reader = PdfReader(file)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            output = os.path.join(UPLOAD_FOLDER, 'signed.pdf')
            with open(output, 'wb') as f:
                writer.write(f)
            return send_file(output, as_attachment=True)
        
        elif tool == 'pdf-validate':
            from pypdf import PdfReader
            file = request.files.get('file')
            try:
                reader = PdfReader(file)
                return jsonify({'valid': True, 'pages': len(reader.pages)})
            except:
                return jsonify({'valid': False})
        
        elif tool == 'pdf-to-word':
            file = request.files.get('file')
            adobe_client_id = os.getenv('ADOBE_CLIENT_ID')
            adobe_client_secret = os.getenv('ADOBE_CLIENT_SECRET')
            
            if adobe_client_id and adobe_client_secret:
                try:
                    print('Using Adobe PDF Services SDK...')
                    from adobe.pdfservices.operation.pdf_services import PDFServices
                    from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
                    from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
                    from adobe.pdfservices.operation.io.stream_asset import StreamAsset
                    from adobe.pdfservices.operation.pdf_services_job import PDFServicesJob
                    from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_params import ExportPDFParams
                    from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_target_format import ExportPDFTargetFormat
                    from adobe.pdfservices.operation.pdfjobs.result.export_pdf_result import ExportPDFResult
                    
                    temp_input = os.path.join(UPLOAD_FOLDER, 'temp_input.pdf')
                    file.save(temp_input)
                    
                    credentials = ServicePrincipalCredentials(
                        client_id=adobe_client_id,
                        client_secret=adobe_client_secret
                    )
                    
                    pdf_services = PDFServices(credentials=credentials)
                    
                    with open(temp_input, 'rb') as f:
                        input_stream = f.read()
                    
                    input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.PDF)
                    export_pdf_params = ExportPDFParams(target_format=ExportPDFTargetFormat.DOCX)
                    
                    from adobe.pdfservices.operation.pdfjobs.jobs.export_pdf_job import ExportPDFJob
                    export_pdf_job = ExportPDFJob(input_asset=input_asset, export_pdf_params=export_pdf_params)
                    
                    location = pdf_services.submit(export_pdf_job)
                    pdf_services_response = pdf_services.get_job_result(location, ExportPDFResult)
                    result_asset = pdf_services_response.get_result().get_asset()
                    stream_asset = pdf_services.get_content(result_asset)
                    
                    output = os.path.join(UPLOAD_FOLDER, 'converted.docx')
                    with open(output, 'wb') as f:
                        f.write(stream_asset.get_input_stream())

                    print('Adobe SDK conversion completed successfully')
                    resp = send_file(output, as_attachment=True)
                    try:
                        resp.headers['X-Conversion-Method'] = 'adobe'
                    except Exception:
                        pass
                    return resp
                except Exception as e:
                    print(f'Adobe SDK failed: {e}, falling back to pdf2docx')
                    import traceback
                    traceback.print_exc()
            
            print('Using pdf2docx...')
            from pdf2docx import Converter
            temp_input = os.path.join(UPLOAD_FOLDER, 'temp_input.pdf')
            output = os.path.join(UPLOAD_FOLDER, 'converted.docx')
            if not os.path.exists(temp_input):
                file.seek(0)
                file.save(temp_input)
            try:
                cv = Converter(temp_input)
                cv.convert(output, start=0, end=None, layout_mode='layout',
                          table_settings={'min_rows_count': 2, 'min_cols_count': 2, 'explicit_borders': True, 'implicit_borders': True},
                          image_settings={'min_width': 10, 'min_height': 10, 'extract_stream': True})
                cv.close()
                print('pdf2docx conversion completed')
                resp = send_file(output, as_attachment=True)
                try:
                    resp.headers['X-Conversion-Method'] = 'pdf2docx'
                except Exception:
                    pass
                return resp
            finally:
                if os.path.exists(temp_input):
                    os.remove(temp_input)
        
        elif tool == 'word-to-pdf':
            file = request.files.get('file')
            adobe_client_id = os.getenv('ADOBE_CLIENT_ID')
            adobe_client_secret = os.getenv('ADOBE_CLIENT_SECRET')
            
            used_method = None
            temp_input = os.path.join(UPLOAD_FOLDER, 'temp_input.docx')
            file.save(temp_input)

            if adobe_client_id and adobe_client_secret:
                try:
                    print('Using Adobe PDF Services SDK for Word to PDF...')
                    from adobe.pdfservices.operation.pdf_services import PDFServices
                    from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
                    from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
                    from adobe.pdfservices.operation.pdfjobs.jobs.create_pdf_job import CreatePDFJob
                    from adobe.pdfservices.operation.pdfjobs.result.create_pdf_result import CreatePDFResult

                    credentials = ServicePrincipalCredentials(
                        client_id=adobe_client_id,
                        client_secret=adobe_client_secret
                    )

                    pdf_services = PDFServices(credentials=credentials)

                    with open(temp_input, 'rb') as f:
                        input_stream = f.read()

                    input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.DOCX)
                    create_pdf_job = CreatePDFJob(input_asset=input_asset)

                    location = pdf_services.submit(create_pdf_job)
                    pdf_services_response = pdf_services.get_job_result(location, CreatePDFResult)
                    result_asset = pdf_services_response.get_result().get_asset()
                    stream_asset = pdf_services.get_content(result_asset)

                    output = os.path.join(UPLOAD_FOLDER, 'converted.pdf')
                    with open(output, 'wb') as f:
                        f.write(stream_asset.get_input_stream())

                    used_method = 'adobe'
                    print('Adobe SDK Word to PDF conversion completed successfully')
                except Exception as e:
                    print(f'Adobe SDK failed: {e}, falling back to Python converters')
                    import traceback
                    traceback.print_exc()

            # If Adobe was not used or failed, try Python-based fallback
            if not used_method:
                try:
                    print('Falling back to mammoth + weasyprint for Word to PDF...')
                    import mammoth
                    from weasyprint import HTML

                    with open(temp_input, 'rb') as f:
                        result = mammoth.convert_to_html(f)
                        html = result.value

                    output = os.path.join(UPLOAD_FOLDER, 'converted.pdf')
                    HTML(string=html).write_pdf(output)
                    used_method = 'mammoth+weasyprint'
                    print('Fallback Word to PDF completed')
                except Exception as e:
                    print(f'Fallback conversion failed: {e}')
                    import traceback
                    traceback.print_exc()
                    if os.path.exists(temp_input):
                        os.remove(temp_input)
                    return jsonify({'error': f'Conversion failed: {str(e)}'}), 500

            # cleanup temp input
            if os.path.exists(temp_input):
                try:
                    os.remove(temp_input)
                except Exception:
                    pass

            resp = send_file(output, as_attachment=True)
            try:
                resp.headers['X-Conversion-Method'] = used_method or 'unknown'
            except Exception:
                pass
            return resp
        
        elif tool == 'pdf-to-excel':
            import fitz
            import pandas as pd
            file = request.files.get('file')
            temp_input = os.path.join(UPLOAD_FOLDER, 'temp_input.pdf')
            file.save(temp_input)
            used_method = None
            try:
                adobe_client_id = os.getenv('ADOBE_CLIENT_ID')
                adobe_client_secret = os.getenv('ADOBE_CLIENT_SECRET')
                if adobe_client_id and adobe_client_secret:
                    try:
                        print('Using Adobe PDF Services SDK for PDF->Excel...')
                        from adobe.pdfservices.operation.pdf_services import PDFServices
                        from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
                        from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
                        from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_params import ExportPDFParams
                        from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_target_format import ExportPDFTargetFormat
                        from adobe.pdfservices.operation.pdfjobs.jobs.export_pdf_job import ExportPDFJob
                        from adobe.pdfservices.operation.pdfjobs.result.export_pdf_result import ExportPDFResult

                        credentials = ServicePrincipalCredentials(client_id=adobe_client_id, client_secret=adobe_client_secret)
                        pdf_services = PDFServices(credentials=credentials)

                        with open(temp_input, 'rb') as f:
                            input_stream = f.read()

                        input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.PDF)
                        export_pdf_params = ExportPDFParams(target_format=ExportPDFTargetFormat.XLSX)
                        export_job = ExportPDFJob(input_asset=input_asset, export_pdf_params=export_pdf_params)

                        location = pdf_services.submit(export_job)
                        pdf_services_response = pdf_services.get_job_result(location, ExportPDFResult)
                        result_asset = pdf_services_response.get_result().get_asset()
                        stream_asset = pdf_services.get_content(result_asset)

                        output = os.path.join(UPLOAD_FOLDER, 'converted.xlsx')
                        with open(output, 'wb') as out_f:
                            out_f.write(stream_asset.get_input_stream())

                        used_method = 'adobe'
                        print('Adobe PDF->Excel conversion completed')
                        resp = send_file(output, as_attachment=True)
                        try:
                            resp.headers['X-Conversion-Method'] = used_method
                        except Exception:
                            pass
                        return resp
                    except Exception as e:
                        print(f'Adobe PDF->Excel failed: {e}, falling back to Python')
                        import traceback
                        traceback.print_exc()

                # Fallback: simple text extraction into Excel
                doc = fitz.open(temp_input)
                all_text = []
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    text = page.get_text()
                    lines = text.split('\n')
                    all_text.extend([line.strip() for line in lines if line.strip()])
                doc.close()
                df = pd.DataFrame({'Content': all_text})
                output = os.path.join(UPLOAD_FOLDER, 'converted.xlsx')
                df.to_excel(output, index=False)
                used_method = 'fitz+pandas'
                resp = send_file(output, as_attachment=True)
                try:
                    resp.headers['X-Conversion-Method'] = used_method
                except Exception:
                    pass
                return resp
            finally:
                if os.path.exists(temp_input):
                    os.remove(temp_input)
        
        elif tool == 'excel-to-pdf':
            import pandas as pd
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            file = request.files.get('file')
            temp_input = os.path.join(UPLOAD_FOLDER, 'temp_input.xlsx')
            file.save(temp_input)
            used_method = None
            try:
                adobe_client_id = os.getenv('ADOBE_CLIENT_ID')
                adobe_client_secret = os.getenv('ADOBE_CLIENT_SECRET')
                if adobe_client_id and adobe_client_secret:
                    try:
                        print('Using Adobe PDF Services SDK for Excel->PDF...')
                        from adobe.pdfservices.operation.pdf_services import PDFServices
                        from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
                        from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
                        from adobe.pdfservices.operation.pdfjobs.jobs.create_pdf_job import CreatePDFJob
                        from adobe.pdfservices.operation.pdfjobs.result.create_pdf_result import CreatePDFResult

                        credentials = ServicePrincipalCredentials(client_id=adobe_client_id, client_secret=adobe_client_secret)
                        pdf_services = PDFServices(credentials=credentials)

                        with open(temp_input, 'rb') as f:
                            input_stream = f.read()

                        input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.XLSX)
                        create_pdf_job = CreatePDFJob(input_asset=input_asset)

                        location = pdf_services.submit(create_pdf_job)
                        pdf_services_response = pdf_services.get_job_result(location, CreatePDFResult)
                        result_asset = pdf_services_response.get_result().get_asset()
                        stream_asset = pdf_services.get_content(result_asset)

                        output = os.path.join(UPLOAD_FOLDER, 'converted.pdf')
                        with open(output, 'wb') as out_f:
                            out_f.write(stream_asset.get_input_stream())

                        used_method = 'adobe'
                        resp = send_file(output, as_attachment=True)
                        try:
                            resp.headers['X-Conversion-Method'] = used_method
                        except Exception:
                            pass
                        return resp
                    except Exception as e:
                        print(f'Adobe Excel->PDF failed: {e}, falling back to Python')
                        import traceback
                        traceback.print_exc()

                # Fallback: simple table render to PDF
                df = pd.read_excel(temp_input)
                output = os.path.join(UPLOAD_FOLDER, 'converted.pdf')
                c = canvas.Canvas(output, pagesize=letter)
                width, height = letter
                y = height - 50
                for index, row in df.iterrows():
                    if y < 50:
                        c.showPage()
                        y = height - 50
                    row_text = ' | '.join([str(val) for val in row.values])
                    c.drawString(50, y, row_text[:100])
                    y -= 20
                c.save()
                used_method = 'pandas+reportlab'
                resp = send_file(output, as_attachment=True)
                try:
                    resp.headers['X-Conversion-Method'] = used_method
                except Exception:
                    pass
                return resp
            finally:
                if os.path.exists(temp_input):
                    os.remove(temp_input)
        
        elif tool == 'pdf-to-powerpoint':
            import fitz
            from pptx import Presentation
            from pptx.util import Inches
            file = request.files.get('file')
            temp_input = os.path.join(UPLOAD_FOLDER, 'temp_input.pdf')
            file.save(temp_input)
            used_method = None
            try:
                adobe_client_id = os.getenv('ADOBE_CLIENT_ID')
                adobe_client_secret = os.getenv('ADOBE_CLIENT_SECRET')
                if adobe_client_id and adobe_client_secret:
                    try:
                        print('Using Adobe PDF Services SDK for PDF->PowerPoint...')
                        from adobe.pdfservices.operation.pdf_services import PDFServices
                        from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
                        from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
                        from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_params import ExportPDFParams
                        from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_target_format import ExportPDFTargetFormat
                        from adobe.pdfservices.operation.pdfjobs.jobs.export_pdf_job import ExportPDFJob
                        from adobe.pdfservices.operation.pdfjobs.result.export_pdf_result import ExportPDFResult

                        credentials = ServicePrincipalCredentials(client_id=adobe_client_id, client_secret=adobe_client_secret)
                        pdf_services = PDFServices(credentials=credentials)

                        with open(temp_input, 'rb') as f:
                            input_stream = f.read()

                        input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.PDF)
                        export_pdf_params = ExportPDFParams(target_format=ExportPDFTargetFormat.PPTX)
                        export_job = ExportPDFJob(input_asset=input_asset, export_pdf_params=export_pdf_params)

                        location = pdf_services.submit(export_job)
                        pdf_services_response = pdf_services.get_job_result(location, ExportPDFResult)
                        result_asset = pdf_services_response.get_result().get_asset()
                        stream_asset = pdf_services.get_content(result_asset)

                        output = os.path.join(UPLOAD_FOLDER, 'converted.pptx')
                        with open(output, 'wb') as out_f:
                            out_f.write(stream_asset.get_input_stream())

                        used_method = 'adobe'
                        resp = send_file(output, as_attachment=True)
                        try:
                            resp.headers['X-Conversion-Method'] = used_method
                        except Exception:
                            pass
                        return resp
                    except Exception as e:
                        print(f'Adobe PDF->PowerPoint failed: {e}, falling back to Python')
                        import traceback
                        traceback.print_exc()

                # Fallback: create PPTX with page text
                doc = fitz.open(temp_input)
                prs = Presentation()
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    text = page.get_text()
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    title.text = f'Page {page_num + 1}'
                    content.text = text[:500]
                doc.close()
                output = os.path.join(UPLOAD_FOLDER, 'converted.pptx')
                prs.save(output)
                used_method = 'fitz+pptx'
                resp = send_file(output, as_attachment=True)
                try:
                    resp.headers['X-Conversion-Method'] = used_method
                except Exception:
                    pass
                return resp
            finally:
                if os.path.exists(temp_input):
                    os.remove(temp_input)
        
        elif tool == 'powerpoint-to-pdf':
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            file = request.files.get('file')
            temp_input = os.path.join(UPLOAD_FOLDER, 'temp_input.pptx')
            file.save(temp_input)
            used_method = None
            try:
                adobe_client_id = os.getenv('ADOBE_CLIENT_ID')
                adobe_client_secret = os.getenv('ADOBE_CLIENT_SECRET')
                if adobe_client_id and adobe_client_secret:
                    try:
                        print('Using Adobe PDF Services SDK for PowerPoint->PDF...')
                        from adobe.pdfservices.operation.pdf_services import PDFServices
                        from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
                        from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
                        from adobe.pdfservices.operation.pdfjobs.jobs.create_pdf_job import CreatePDFJob
                        from adobe.pdfservices.operation.pdfjobs.result.create_pdf_result import CreatePDFResult

                        credentials = ServicePrincipalCredentials(client_id=adobe_client_id, client_secret=adobe_client_secret)
                        pdf_services = PDFServices(credentials=credentials)

                        with open(temp_input, 'rb') as f:
                            input_stream = f.read()

                        input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.PPTX)
                        create_pdf_job = CreatePDFJob(input_asset=input_asset)

                        location = pdf_services.submit(create_pdf_job)
                        pdf_services_response = pdf_services.get_job_result(location, CreatePDFResult)
                        result_asset = pdf_services_response.get_result().get_asset()
                        stream_asset = pdf_services.get_content(result_asset)

                        output = os.path.join(UPLOAD_FOLDER, 'converted.pdf')
                        with open(output, 'wb') as out_f:
                            out_f.write(stream_asset.get_input_stream())

                        used_method = 'adobe'
                        resp = send_file(output, as_attachment=True)
                        try:
                            resp.headers['X-Conversion-Method'] = used_method
                        except Exception:
                            pass
                        return resp
                    except Exception as e:
                        print(f'Adobe PowerPoint->PDF failed: {e}, falling back to Python')
                        import traceback
                        traceback.print_exc()

                # Fallback: simple text render from PPTX to PDF
                prs = Presentation(temp_input)
                output = os.path.join(UPLOAD_FOLDER, 'converted.pdf')
                c = canvas.Canvas(output, pagesize=letter)
                width, height = letter
                for slide_num, slide in enumerate(prs.slides):
                    if slide_num > 0:
                        c.showPage()
                    y = height - 50
                    c.drawString(50, y, f'Slide {slide_num + 1}')
                    y -= 30
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            if y < 50:
                                c.showPage()
                                y = height - 50
                            c.drawString(50, y, shape.text[:80])
                            y -= 20
                c.save()
                used_method = 'pptx+reportlab'
                resp = send_file(output, as_attachment=True)
                try:
                    resp.headers['X-Conversion-Method'] = used_method
                except Exception:
                    pass
                return resp
            finally:
                if os.path.exists(temp_input):
                    os.remove(temp_input)
        
        return jsonify({'error': 'Tool not implemented'}), 400
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5001)
